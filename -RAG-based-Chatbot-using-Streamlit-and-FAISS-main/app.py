import faiss
import os
import streamlit as st
from io import BytesIO
from docx import Document
import numpy as np
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.chains import RetrievalQA
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.docstore.in_memory import InMemoryDocstore
from langchain_huggingface import HuggingFaceEndpoint

# Set Hugging Face API key
huggingface_api_key = st.secrets["huggingface_api_key"]
os.environ['HUGGINGFACEHUB_API_TOKEN'] = huggingface_api_key

# ---------------------------- TEXT EXTRACTION -------------------------------- #
def extract_text_from_docx(file):
    doc = Document(file)
    full_text = [para.text for para in doc.paragraphs]
    return '\n'.join(full_text)

def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(input_type, input_data):
    file = input_data if isinstance(input_data, BytesIO) else BytesIO(input_data.read())

    if input_type == "PDF":
        pdf_reader = PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            if page.extract_text():
                text += page.extract_text()
    elif input_type == "DOCX":
        text = extract_text_from_docx(file)
    elif input_type == "TXT":
        text = file.read().decode("utf-8")
    elif input_type == "PPTX":
        text = extract_text_from_pptx(file)
    else:
        raise ValueError("Unsupported file type")
    return text

# ---------------------------- VECTOR STORE -------------------------------- #
def process_input(input_type, input_data, full_text):
    # Split into chunks
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    texts = text_splitter.split_text(full_text)

    # Embeddings
    model_name = "sentence-transformers/all-mpnet-base-v2"
    model_kwargs = {'device': 'cpu'}
    encode_kwargs = {'normalize_embeddings': False}
    hf_embeddings = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs=model_kwargs,
        encode_kwargs=encode_kwargs
    )

    # FAISS index
    sample_embedding = np.array(hf_embeddings.embed_query("sample text"))
    dimension = sample_embedding.shape[0]
    index = faiss.IndexFlatL2(dimension)

    vector_store = FAISS(
        embedding_function=hf_embeddings.embed_query,
        index=index,
        docstore=InMemoryDocstore(),
        index_to_docstore_id={},
    )
    vector_store.add_texts(texts)
    return vector_store

# ---------------------------- QA FUNCTION -------------------------------- #
def answer_question(vectorstore, query):
    llm = HuggingFaceEndpoint(
        repo_id='HuggingFaceH4/zephyr-7b-beta',
        token=huggingface_api_key,
        task="text-generation",
        temperature=0.6
    )
    qa = RetrievalQA.from_chain_type(llm=llm, retriever=vectorstore.as_retriever())
    answer = qa({"query": query})
    return answer["result"]

# ---------------------------- MAIN APP -------------------------------- #
def main():
    st.set_page_config(page_title="RAG Q&A App", layout="wide")
    st.title("🧠 RAG-based Document Q&A")

    # Initialize session state
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "vectorstore" not in st.session_state:
        st.session_state.vectorstore = None

    # --- CSS Styling ---
    st.markdown("""
        <style>
            .chat-bubble {
                padding: 15px;
                margin: 10px 0;
                border-radius: 10px;
                font-family: sans-serif;
                font-size: 15px;
                white-space: pre-wrap;
            }
            .user-msg {
                background-color: #DCF8C6;
                text-align: right;
                border: 1px solid #ccc;
            }
            .bot-msg {
                background-color: #F1F0F0;
                text-align: left;
                border: 1px solid #ccc;
            }
            .preview-box {
                background-color: #f9f9f9;
                padding: 20px;
                border-radius: 10px;
                height: 400px;
                overflow-y: auto;
                font-family: monospace;
                font-size: 14px;
                white-space: pre-wrap;
                border: 1px solid #ddd;
            }
        </style>
    """, unsafe_allow_html=True)

    # --- Sidebar for Upload ---
    with st.sidebar:
        st.header("📄 Upload & Process")
        input_type = st.selectbox("Input Type", ["PDF", "DOCX", "TXT", "PPTX"])
        file_types = {
            "PDF": ["pdf"],
            "DOCX": ["docx"],
            "TXT": ["txt"],
            "PPTX": ["pptx"]
        }
        input_data = st.file_uploader(f"Upload a {input_type} file", type=file_types[input_type])

        if st.button("🧹 Clear Chat History"):
            st.session_state.chat_history = []
            st.success("Chat history cleared.")

    # --- Document Preview and Processing ---
    col1, col2 = st.columns([1, 2])

    if input_data is not None:
        extracted_text = extract_text(input_type, input_data)
        with col1:
            st.markdown("### 📄 File Preview")
            st.markdown(f"<div class='preview-box'>{extracted_text}</div>", unsafe_allow_html=True)

        with col2:
            if st.button("📥 Process Document"):
                with st.spinner("Indexing document..."):
                    vectorstore = process_input(input_type, input_data, extracted_text)
                    st.session_state.vectorstore = vectorstore
                    st.success("Document processed and indexed successfully!")

    # --- Q&A Section ---
    if st.session_state.vectorstore:
        st.markdown("### 💬 Ask your Question")
        query = st.text_input("Type your question here...")

        if st.button("🔍 Submit"):
            if query.strip():
                answer = answer_question(st.session_state.vectorstore, query)
                st.session_state.chat_history.append(("user", query))
                st.session_state.chat_history.append(("bot", answer))
            else:
                st.warning("Please enter a question.")

        # Display chat
        for speaker, msg in st.session_state.chat_history:
            css_class = "user-msg" if speaker == "user" else "bot-msg"
            st.markdown(f"<div class='chat-bubble {css_class}'>{msg}</div>", unsafe_allow_html=True)

if __name__ == "__main__":
    main() 