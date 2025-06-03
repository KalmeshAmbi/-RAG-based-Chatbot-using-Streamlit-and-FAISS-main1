import streamlit as st
import faiss
import os
from io import BytesIO
from docx import Document as DocxDocument
import numpy as np
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.llms import HuggingFaceEndpoint
from langchain.schema import Document
from pptx import Presentation
from PyPDF2 import PdfReader

# ---------------- CONFIG & SETUP ----------------

huggingface_api_key = st.secrets["huggingface_api_key"]
os.environ['HUGGINGFACEHUB_API_TOKEN'] = huggingface_api_key

# ----------------- HELPERS ----------------------

def truncate_documents(docs, max_chars=8000):
    total_text = ""
    result_docs = []
    for doc in docs:
        if len(total_text) + len(doc.page_content) <= max_chars:
            result_docs.append(doc)
            total_text += doc.page_content
        else:
            break
    return result_docs

def extract_text_from_file(input_type, input_data):
    text = ""
    try:
        if input_type == "PDF":
            pdf_reader = PdfReader(input_data)
            for page in pdf_reader.pages:
                if page.extract_text():
                    text += page.extract_text()
        elif input_type == "Text":
            input_data.seek(0)  # Reset file pointer to the beginning
            try:
                text = input_data.read().decode("utf-8")  # Try UTF-8 first
            except UnicodeDecodeError:
                text = input_data.read().decode("ISO-8859-1")  # Fallback encoding
        elif input_type == "PPT":
            presentation = Presentation(input_data)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif input_type == "DOCX":
            doc = DocxDocument(input_data)
            for para in doc.paragraphs:
                text += para.text + "\n"
    except Exception as e:
        st.error(f"❌ Error while extracting text: {e}")
    return text

def process_input(input_type, input_data):
    raw_text = extract_text_from_file(input_type, input_data)

    if not raw_text.strip():
        st.warning("Uploaded document is empty or could not be read.")
        return None

    # Split the text into manageable chunks
    text_splitter = CharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    texts = text_splitter.split_text(raw_text)

    if not texts:
        st.warning("No text could be extracted for processing. Please check your document.")
        return None

    docs = [Document(page_content=chunk) for chunk in texts]

    # Embedding model setup
    model_name = "sentence-transformers/all-mpnet-base-v2"
    hf_embeddings = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs={'device': 'cpu'},
        encode_kwargs={'normalize_embeddings': False}
    )

    # Create FAISS vector store
    vector_store = FAISS.from_documents(docs, hf_embeddings)
    return vector_store

def answer_question(vectorstore, query):
    if not query.strip():
        return {"result": "⚠️ Please enter a valid question."}

    # LLM endpoint setup
    llm = HuggingFaceEndpoint(
        repo_id="google/flan-t5-large",
        task="text2text-generation",
        huggingfacehub_api_token=huggingface_api_key,
        temperature=0.6
    )

    # Retrieve relevant documents
    retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
    docs = retriever.get_relevant_documents(query)

    if not docs:
        return {"result": "⚠️ No relevant content found to answer your question."}

    # Truncate documents to fit the context
    truncated_docs = truncate_documents(docs)
    context = "\n\n".join([doc.page_content for doc in truncated_docs])
    prompt = f"""Answer the following question based on the context:\n\nContext: {context}\n\nQuestion: {query}"""

   
        return {"result": response}
    

def format_document_content(input_type, input_data):
    document_content = extract_text_from_file(input_type, input_data)
    if input_type in ["PPT", "DOCX"]:
        lines = document_content.split("\n")
        document_content = "".join(f"<p>{line}</p>" for line in lines if line.strip())
    return document_content

# ---------------- MAIN APP ----------------------

def main():
    st.set_page_config(page_title="RAG Q&A App", layout="wide")
    st.title("🧠 RAG Q&A Application")

    if "chat_history" not in st.session_state:
        st.session_state["chat_history"] = []

    st.markdown("""
        <style>
            .document-preview {
                max-height: 500px;
                overflow-y: scroll;
                background-color: #ffffff;
                padding: 15px;
                border-radius: 8px;
                border: 1px solid #ccc;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
                margin-bottom: 20px;
            }
            .green-button > button {
                background-color: #4CAF50 !important;
                color: white !important;
                border: none;
                padding: 10px 20px;
                border-radius: 5px;
            }
        </style>
    """, unsafe_allow_html=True)

    col1, col2 = st.columns([2, 3])

    file_types = {"PDF": "pdf", "Text": "txt", "PPT": "pptx", "DOCX": "docx"}

    with col1:
        st.subheader("Upload Your Document")
        input_type = st.selectbox("Choose Input Type", list(file_types.keys()))
        input_data = None
        document_content = ""

        if input_type:
            file_extension = file_types.get(input_type)
            input_data = st.file_uploader(f"Upload a {input_type} file", type=[file_extension])

        if input_data:
            document_content = format_document_content(input_type, input_data)
            st.subheader("Document Content")
            st.markdown(f'<div class="document-preview">{document_content}</div>', unsafe_allow_html=True)

    with col2:
        st.subheader("Ask Questions Based on the Document")

        if input_data and st.button("🔄 Process Document"):
            vectorstore = process_input(input_type, input_data)
            if vectorstore:
                st.session_state["vectorstore"] = vectorstore
                st.success("✅ Document processed successfully.")
            else:
                st.error("❌ Document processing failed. Check if the file has readable content.")

        if "vectorstore" in st.session_state:
            query = st.text_input("Ask your question:")
            st.markdown('<div class="green-button">', unsafe_allow_html=True)
            if st.button("✅ Submit Question"):
                answer = answer_question(st.session_state["vectorstore"], query)
                st.session_state["chat_history"].append({"question": query, "answer": answer["result"]})
                st.session_state["latest_answer"] = answer["result"]
            st.markdown('</div>', unsafe_allow_html=True)

        if "latest_answer" in st.session_state:
            st.markdown(f"""
                <div style='background-color:#e6f7ff; padding:20px; border-left:6px solid #2196F3;
                            border-radius:8px; margin-top:20px; box-shadow: 0px 1px 3px rgba(0,0,0,0.1);'>
                    <h4 style='color:#0b5394;'>Answer:</h4>
                    <p style='font-size:16px; color:#333;'>{st.session_state["latest_answer"]}</p>
                </div>
            """, unsafe_allow_html=True)

    st.subheader("Chat History")
    if st.button("🗑️ Clear Chat History", use_container_width=True, type="primary"):
        st.session_state["chat_history"] = []
        st.session_state["latest_answer"] = ""
        st.success("🧹 Chat history cleared.")

    if st.session_state["chat_history"]:
        for chat in reversed(st.session_state["chat_history"]):
            st.markdown(f"""
            <div style='background-color:#f9f9f9; padding:15px; border-radius:10px; margin-bottom:10px;
                        box-shadow: 0px 1px 3px rgba(0,0,0,0.1);'>
                <strong style='color:#333;'>You:</strong> {chat['question']}<br>
                <strong style='color:#4CAF50;'>Answer:</strong> {chat['answer']}
            </div>
            """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
